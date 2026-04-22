"""
Generate Final Review PPT for Team 22 - IBS Project
Uses the same template style as Team23_in_MovieTemplate.pptx
"""
import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from copy import deepcopy

# ──────────────────────────────────────────────
# CONFIGURATION
# ──────────────────────────────────────────────
TEMPLATE_PATH = r'c:\College\Me\College\4th Sem\IBS\GNN_Project\Team23_in_MovieTemplate.pptx'
OUTPUT_PATH = r'c:\College\Me\College\4th Sem\IBS\GNN_Project\Team-22_Final_Review.pptx'
FONT_NAME = 'Cambria Math'

# Project images
PIPELINE_IMG = r'C:\Users\Achintya\.gemini\antigravity\brain\624357d7-9b96-4147-b599-909d193b1895\pipeline_architecture_1774629672232.png'
MODEL_IMG = r'C:\Users\Achintya\.gemini\antigravity\brain\624357d7-9b96-4147-b599-909d193b1895\model_architecture_1774629749685.png'
GRAPH_IMG = r'C:\Users\Achintya\.gemini\antigravity\brain\624357d7-9b96-4147-b599-909d193b1895\graph_construction_1774629824978.png'
TRAINING_CURVE_GCN = r'c:\College\Me\College\4th Sem\IBS\GNN_Project\training_curve.png'
TRAINING_CURVE_GIN = r'c:\College\Me\College\4th Sem\IBS\GNN_Project\gin_training_curve.png'
SCATTER_GCN = r'c:\College\Me\College\4th Sem\IBS\GNN_Project\prediction_scatter.png'
SCATTER_GIN = r'c:\College\Me\College\4th Sem\IBS\GNN_Project\gin_prediction_scatter.png'

# Template images extracted
LOGO_IMG = r'c:\College\Me\College\4th Sem\IBS\GNN_Project\template_images\slide1_Image_0.png'
FOOTER_IMG = r'c:\College\Me\College\4th Sem\IBS\GNN_Project\template_images\slide1_Image_1.png'

# Colors from template
MAROON = RGBColor(0x8B, 0x1A, 0x3B)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_BG = RGBColor(0x8B, 0x1A, 0x3B)

# Sizes from template
SLIDE_W = 9144000
SLIDE_H = 5143500
HEADER_H = 566928
LOGO_SIZE = 502920
LOGO_LEFT = 54864
LOGO_TOP = 36576
TITLE_LEFT = 658368
TITLE_TOP = 64008
TITLE_W = 8229600
TITLE_H = 457200
FOOTER_LEFT = 6583680
FOOTER_TOP = 4434840
FOOTER_W = 2377440
FOOTER_H = 594360

# Font sizes from template
TITLE_FONT = Pt(22)  # 279400 EMU ≈ 22pt
SECTION_HEADER_FONT = Pt(13)  # 165100 EMU ≈ 13pt
BODY_FONT = Pt(12)  # 152400 EMU ≈ 12pt
BODY_SMALL_FONT = Pt(11)  # 139700 EMU ≈ 11pt
THANK_YOU_FONT = Pt(44)  # 558800 EMU ≈ 44pt
SUBTITLE_FONT = Pt(20)  # 254000 EMU ≈ 20pt
SUBTITLE_SMALL_FONT = Pt(14)  # 177800 EMU ≈ 14pt

# Content body area
BODY_LEFT = 182880
BODY_TOP = 590000
BODY_W = 8705088
BODY_H = 3700000


class SlideBuilder:
    """Helper class to build slides matching the Team23 template style."""
    
    def __init__(self, prs):
        self.prs = prs
        self.layout = prs.slide_layouts[0]  # DEFAULT layout
    
    def add_slide(self):
        slide = self.prs.slides.add_slide(self.layout)
        # Remove all default placeholders
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        return slide
    
    def add_header_bar(self, slide, title_text, slide_num=None):
        """Add the maroon header bar with logo and title."""
        # Header bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = HEADER_BG
        bar.line.fill.background()
        bar.name = 'Shape 0'
        
        # Logo
        if os.path.exists(LOGO_IMG):
            logo = slide.shapes.add_picture(LOGO_IMG, LOGO_LEFT, LOGO_TOP, LOGO_SIZE, LOGO_SIZE)
            logo.name = 'Image 0'
        
        # Title text box
        title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
        title_box.fill.background()
        title_box.line.fill.background()
        title_box.name = 'Text 1'
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        if slide_num:
            run.text = f"{slide_num}  {title_text}"
        else:
            run.text = title_text
        run.font.name = FONT_NAME
        run.font.size = TITLE_FONT
        run.font.bold = True
        run.font.color.rgb = WHITE
    
    def add_footer(self, slide):
        """Add the footer image (Amrita logo)."""
        if os.path.exists(FOOTER_IMG):
            footer = slide.shapes.add_picture(FOOTER_IMG, FOOTER_LEFT, FOOTER_TOP, FOOTER_W, FOOTER_H)
            footer.name = 'Image 1'
    
    def add_text_box(self, slide, left, top, width, height, name='Text'):
        """Add a transparent text box."""
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.background()
        box.line.fill.background()
        box.name = name
        box.text_frame.word_wrap = True
        return box.text_frame
    
    def add_section(self, tf, title):
        """Add a section header (maroon bold)."""
        if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = title
        run.font.name = FONT_NAME
        run.font.size = SECTION_HEADER_FONT
        run.font.bold = True
        run.font.color.rgb = MAROON
        return p
    
    def add_bullet(self, tf, text, font_size=None, bold=False, color=None):
        """Add a bullet point."""
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = font_size or BODY_SMALL_FONT
        run.font.bold = bold
        run.font.color.rgb = color or DARK_TEXT
        return p
    
    def add_blank_line(self, tf):
        """Add an empty paragraph for spacing."""
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = ""
        run.font.size = Pt(6)
        return p


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    
    # Add a blank layout
    # We need to create slides from scratch since we're not using the template's slides
    # Let's use the template file but clear its slides
    prs = Presentation(TEMPLATE_PATH)
    
    # Delete all existing slides
    slide_ids = list(prs.slides._sldIdLst)
    for sldId in slide_ids:
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)
    
    sb = SlideBuilder(prs)
    
    # ════════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    
    # Full maroon background for title slide
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER_BG
    bar.line.fill.background()
    
    # Logo top-left
    if os.path.exists(LOGO_IMG):
        slide.shapes.add_picture(LOGO_IMG, LOGO_LEFT, LOGO_TOP, LOGO_SIZE, LOGO_SIZE)
    
    # Team info top (Lishanth first — matching original Review 1 & 2 PPTs)
    tf = sb.add_text_box(slide, Emu(914400), Emu(700000), Emu(7315200), Emu(500000), 'Text 1')
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "CH.SC.U4AIE24083 \u2013 Lishanth SS"
    run.font.name = FONT_NAME
    run.font.size = Pt(14)
    run.font.bold = False
    run.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "CH.SC.U4AIE24003 \u2013 Achintya TR"
    run.font.name = FONT_NAME
    run.font.size = Pt(14)
    run.font.bold = False
    run.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    
    # Main title
    tf = sb.add_text_box(slide, Emu(914400), Emu(1500000), Emu(7315200), Emu(1800000), 'Text 2')
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Structure-Guided Graph Neural Network for Predicting Protein Melting Temperature from Sequence and 3D Features"
    run.font.name = FONT_NAME
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    
    # Subtitle info
    tf = sb.add_text_box(slide, Emu(914400), Emu(3500000), Emu(7315200), Emu(600000), 'Text 3')
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "TEAM-22  |  IBS-2 Project  |  Semester 4"
    run.font.name = FONT_NAME
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Department of Computer Science & Engineering \u2013 AI"
    run.font.name = FONT_NAME
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Faculty Mentor: Dr. I. R. Oviya"
    run.font.name = FONT_NAME
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    
    # Footer
    if os.path.exists(FOOTER_IMG):
        slide.shapes.add_picture(FOOTER_IMG, FOOTER_LEFT, FOOTER_TOP, FOOTER_W, FOOTER_H)
    
    # ════════════════════════════════════════════
    # SLIDE 2: Problem Statement (matches Review 1 Slide 2)
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "PROBLEM STATEMENT", "01")
    sb.add_footer(slide)
    
    tf = sb.add_text_box(slide, BODY_LEFT, BODY_TOP, BODY_W, BODY_H, 'Text 2')
    sb.add_section(tf, "The Challenge")
    sb.add_bullet(tf, "Protein engineers design novel sequences for therapeutics and industrial applications, yet ensuring sufficient thermal stability (high melting temperature, Tm) remains a critical bottleneck.")
    sb.add_bullet(tf, "Experimental techniques such as differential scanning calorimetry are costly, time-consuming, and unsuitable for screening vast sequence spaces.")
    sb.add_bullet(tf, "Existing computational methods either focus on relative stability changes (\u0394\u0394G) for point mutations or rely solely on sequence information, limiting generalization to diverse folds and de novo designs.")
    sb.add_blank_line(tf)
    sb.add_section(tf, "Our Solution")
    sb.add_bullet(tf, "This project addresses the challenge by framing Tm prediction as a regression task using both sequence features and 3D protein graphs.")
    sb.add_bullet(tf, "By integrating protein-language-model (ESM-2) embeddings and graph neural networks on residue-contact graphs from AlphaFold structures, the goal is to build a robust system that accurately predicts absolute Tm across diverse proteins.")

    # ════════════════════════════════════════════
    # SLIDE 3: Dataset Overview (matches Review 2 Slides 2-3)
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "DATASET OVERVIEW", "02")
    sb.add_footer(slide)
    
    # Left column: Dataset source
    tf = sb.add_text_box(slide, BODY_LEFT, BODY_TOP, 4300000, 3750000, 'TextLeft')
    sb.add_section(tf, "DATA SOURCE")
    sb.add_bullet(tf, "Dataset from TemStaPro study (Bioinformatics, 2024)")
    sb.add_bullet(tf, "Contains experimentally reported melting temperature (Tm in \u00b0C)")
    sb.add_bullet(tf, "Includes metadata: Protein name, Source organism, pH, Method, PubMed ID")
    sb.add_bullet(tf, "Some records also contain stability-related values like \u0394G")
    sb.add_bullet(tf, "Total: 943,605 protein sequences with Tm labels")
    sb.add_blank_line(tf)
    sb.add_section(tf, "DATASET CHARACTERISTICS")
    sb.add_bullet(tf, "Proteins from diverse organisms \u2192 supports generalization")
    sb.add_bullet(tf, "Multiple entries per protein under different conditions")
    sb.add_bullet(tf, "Missing values present \u2192 preprocessing required")
    
    # Right column: Why this dataset
    tf = sb.add_text_box(slide, 4660000, BODY_TOP, 4300000, 3750000, 'TextRight')
    sb.add_section(tf, "WHY THIS DATASET")
    sb.add_bullet(tf, "Suitable for ML regression models for thermostability")
    sb.add_bullet(tf, "Enables using both experimental conditions and sequence embeddings")
    sb.add_bullet(tf, "Supports baseline models and GNN-based approaches")
    sb.add_blank_line(tf)
    sb.add_section(tf, "DATA SPLITS")
    sb.add_bullet(tf, "80% Training / 20% Testing split")
    sb.add_bullet(tf, "Development subset: 500 samples (rapid prototyping)")
    sb.add_bullet(tf, "Intermediate: 5K \u2192 20K subsets")
    sb.add_bullet(tf, "Final training: 100,000 samples")

    # ════════════════════════════════════════════
    # SLIDE 4: Introduction  
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "INTRODUCTION", "03")
    sb.add_footer(slide)
    
    # Left column
    tf = sb.add_text_box(slide, BODY_LEFT, BODY_TOP, 4300000, 3750000, 'TextLeft')
    sb.add_section(tf, "THE PROBLEM")
    sb.add_bullet(tf, "Protein thermal stability (Tm) is critical for drug design & industrial enzymes")
    sb.add_bullet(tf, "Experimental Tm measurement is slow, expensive, and low-throughput")
    sb.add_bullet(tf, "Sequence-only models miss critical 3D structural stability features")
    sb.add_bullet(tf, "Most existing GNNs predict \u0394\u0394G for mutations, not absolute Tm")
    sb.add_blank_line(tf)
    sb.add_section(tf, "KEY HIGHLIGHTS")
    sb.add_bullet(tf, "Dataset: 943,605 proteins from TemStaPro (Bioinformatics, 2024)")
    sb.add_bullet(tf, "Features: ESM-2 protein language model embeddings (320-dim)")
    sb.add_bullet(tf, "Structures: AlphaFold Database predicted structures")
    sb.add_bullet(tf, "Two GNN architectures compared: GCN vs GIN")
    sb.add_bullet(tf, "Trained on 100,000 protein structures locally")
    
    # Right column
    tf = sb.add_text_box(slide, 4660000, BODY_TOP, 4300000, 3750000, 'TextRight')
    sb.add_section(tf, "OUR APPROACH")
    sb.add_bullet(tf, "Represent proteins as residue-level graphs from 3D structures")
    sb.add_bullet(tf, "Nodes = amino acids with ESM-2 deep embeddings")
    sb.add_bullet(tf, "Edges = spatial contacts (C\u03b1 distance < 8 \u00c5)")
    sb.add_bullet(tf, "GNN learns structural patterns predictive of stability")
    sb.add_bullet(tf, "Direct Tm regression \u2013 not mutation-based \u0394\u0394G")

    # ════════════════════════════════════════════
    # SLIDE 5: Objective
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "OBJECTIVE", "04")
    sb.add_footer(slide)
    
    tf = sb.add_text_box(slide, BODY_LEFT, BODY_TOP, BODY_W, BODY_H, 'Text 2')
    
    sb.add_section(tf, "Primary Goal")
    sb.add_bullet(tf, "Develop a structure-guided Graph Neural Network to predict the absolute melting temperature (Tm) of proteins using both amino acid sequence and 3D structural information.")
    sb.add_blank_line(tf)
    sb.add_section(tf, "Innovation")
    sb.add_bullet(tf, "First integration of pre-trained Protein Language Model (ESM-2) embeddings with structure-based GNN for direct Tm regression on the TemStaPro large-scale dataset (943K proteins).")
    sb.add_blank_line(tf)
    sb.add_section(tf, "Comparative Study")
    sb.add_bullet(tf, "Systematically compare GCN (Graph Convolutional Network) vs GIN (Graph Isomorphism Network) architectures to identify optimal graph learning strategy for Tm prediction.")
    sb.add_blank_line(tf)
    sb.add_section(tf, "Practical Impact")
    sb.add_bullet(tf, "Enable rapid computational screening of protein thermostability for drug design, industrial enzyme engineering, and de novo protein design applications.")
    
    # ════════════════════════════════════════════
    # SLIDE 6: Methodology
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "METHODOLOGY", "05")
    sb.add_footer(slide)
    
    # Three columns like template slide 4
    # Column 1: Data Preparation
    tf = sb.add_text_box(slide, BODY_LEFT, BODY_TOP, 2880000, 3700000, 'TextCol0')
    sb.add_section(tf, "DATA PREPARATION")
    sb.add_bullet(tf, "943,605 sequences from TemStaPro FASTA dataset", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Parsed headers \u2192 UniProt ID, Sequence, Tm", BODY_SMALL_FONT)
    sb.add_bullet(tf, "80/20 Train/Test split", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Subsets: 500 (dev) \u2192 5K \u2192 20K \u2192 100K", BODY_SMALL_FONT)
    
    # Column 2: Structure & Features
    tf = sb.add_text_box(slide, 3120000, BODY_TOP, 2880000, 3700000, 'TextCol1')
    sb.add_section(tf, "STRUCTURE & FEATURES")
    sb.add_bullet(tf, "3D structures from AlphaFold Database API", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Automated download with version resolution", BODY_SMALL_FONT)
    sb.add_bullet(tf, "ESM-2 embeddings (esm2_t6_8M, dim=320)", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Graph: Nodes = residues, Edges = C\u03b1 < 8\u00c5", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Fallback: linear graph if PDB missing", BODY_SMALL_FONT)
    
    # Column 3: Model & Training
    tf = sb.add_text_box(slide, 6057120, BODY_TOP, 2880000, 3700000, 'TextCol2')
    sb.add_section(tf, "MODEL & TRAINING")
    sb.add_bullet(tf, "GCN: 3\u00d7 GCNConv + BatchNorm + ReLU", BODY_SMALL_FONT)
    sb.add_bullet(tf, "GIN: 3\u00d7 GINConv + MLP + BatchNorm", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Global Mean Pooling \u2192 MLP regression", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Dropout (p=0.5) + L2 regularization", BODY_SMALL_FONT)
    sb.add_blank_line(tf)
    sb.add_section(tf, "TRAINING CONFIG")
    sb.add_bullet(tf, "Optimizer: Adam (lr=0.001)", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Loss: MSE | Epochs: 30", BODY_SMALL_FONT)
    sb.add_bullet(tf, "Early stopping on test loss", BODY_SMALL_FONT)
    sb.add_bullet(tf, "GPU: NVIDIA RTX 4050 (CUDA)", BODY_SMALL_FONT)
    
    # ════════════════════════════════════════════
    # SLIDE 7: Pipeline Architecture Diagram  
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "PIPELINE ARCHITECTURE", "06")
    sb.add_footer(slide)
    
    if os.path.exists(PIPELINE_IMG):
        # Center the image
        img_w = 8400000
        img_h = 3600000
        left = (SLIDE_W - img_w) // 2
        top = HEADER_H + 150000
        slide.shapes.add_picture(PIPELINE_IMG, left, top, img_w, img_h)
    
    # ════════════════════════════════════════════
    # SLIDE 8: Graph Construction + Model Architecture
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "GRAPH CONSTRUCTION & MODEL", "07")
    sb.add_footer(slide)
    
    # Left: Graph construction image
    if os.path.exists(GRAPH_IMG):
        slide.shapes.add_picture(GRAPH_IMG, BODY_LEFT, HEADER_H + 100000, 4200000, 3200000)
    
    # Right: Model architecture image 
    if os.path.exists(MODEL_IMG):
        slide.shapes.add_picture(MODEL_IMG, 4600000, HEADER_H + 100000, 4200000, 3200000)

    # ════════════════════════════════════════════
    # SLIDE 9: Results \u2013 Training Curves
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "RESULTS \u2013 TRAINING CURVES", "08")
    sb.add_footer(slide)
    
    # Left: GCN training curve
    if os.path.exists(TRAINING_CURVE_GCN):
        slide.shapes.add_picture(TRAINING_CURVE_GCN, BODY_LEFT, HEADER_H + 100000, 4200000, 3200000)
    
    # Right: GIN training curve
    if os.path.exists(TRAINING_CURVE_GIN):
        slide.shapes.add_picture(TRAINING_CURVE_GIN, 4600000, HEADER_H + 100000, 4200000, 3200000)
    
    # Labels
    tf = sb.add_text_box(slide, BODY_LEFT, HEADER_H + 3400000, 4200000, 300000, 'LabelLeft')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "GCN Training Curve (100K Dataset)"
    run.font.name = FONT_NAME
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = MAROON
    
    tf = sb.add_text_box(slide, 4600000, HEADER_H + 3400000, 4200000, 300000, 'LabelRight')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "GIN Training Curve (100K Dataset)"
    run.font.name = FONT_NAME
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = MAROON

    # ════════════════════════════════════════════
    # SLIDE 10: Results \u2013 Scatter Plots
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "RESULTS \u2013 PREDICTION ACCURACY", "09")
    sb.add_footer(slide)
    
    # Left: GCN scatter
    if os.path.exists(SCATTER_GCN):
        slide.shapes.add_picture(SCATTER_GCN, BODY_LEFT, HEADER_H + 100000, 4200000, 3200000)
    
    # Right: GIN scatter
    if os.path.exists(SCATTER_GIN):
        slide.shapes.add_picture(SCATTER_GIN, 4600000, HEADER_H + 100000, 4200000, 3200000)
    
    # Labels
    tf = sb.add_text_box(slide, BODY_LEFT, HEADER_H + 3400000, 4200000, 300000, 'LabelLeft')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "GCN: Pearson r = 0.720"
    run.font.name = FONT_NAME
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = MAROON
    
    tf = sb.add_text_box(slide, 4600000, HEADER_H + 3400000, 4200000, 300000, 'LabelRight')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "GIN: Pearson r = 0.744 (Best)"
    run.font.name = FONT_NAME
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = MAROON

    # ════════════════════════════════════════════
    # SLIDE 11: Model Comparison Table + Discussion
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "MODEL COMPARISON", "10")
    sb.add_footer(slide)
    
    # Left side: Results table using text boxes (tables are complex with python-pptx)
    tf = sb.add_text_box(slide, BODY_LEFT, BODY_TOP, 4300000, 3700000, 'TextLeft')
    sb.add_section(tf, "QUANTITATIVE RESULTS (100K DATASET)")
    sb.add_blank_line(tf)
    sb.add_bullet(tf, "                    GCN          GIN", Pt(10), bold=True, color=MAROON)
    sb.add_bullet(tf, "Pearson r:       0.720        0.744", Pt(11))
    sb.add_bullet(tf, "Test MSE:        101.40       95.12", Pt(11))
    sb.add_bullet(tf, "Test RMSE:      10.07         9.75", Pt(11))
    sb.add_bullet(tf, "Test MAE:        7.54         7.23", Pt(11))
    sb.add_bullet(tf, "Test R\u00b2:           0.51         0.55", Pt(11))
    sb.add_blank_line(tf)
    sb.add_section(tf, "SCALING PROGRESSION")
    sb.add_bullet(tf, "500 samples \u2192 MSE 101 (baseline)", Pt(11))
    sb.add_bullet(tf, "20K samples \u2192 Pearson 0.66 (early stop)", Pt(11))
    sb.add_bullet(tf, "100K samples \u2192 Pearson 0.744 (best)", Pt(11))
    
    # Right side: Discussion
    tf = sb.add_text_box(slide, 4660000, BODY_TOP, 4300000, 3700000, 'TextRight')
    sb.add_section(tf, "KEY FINDINGS")
    sb.add_bullet(tf, "GIN outperforms GCN by ~3.3% in Pearson correlation")
    sb.add_bullet(tf, "GIN's injective aggregation better captures structural motifs")
    sb.add_bullet(tf, "ESM-2 embeddings provide far richer features than one-hot encoding")
    sb.add_bullet(tf, "3D structural contacts capture long-range stability interactions")
    sb.add_blank_line(tf)
    sb.add_section(tf, "COMPARISON WITH LITERATURE")
    sb.add_bullet(tf, "Outperforms sequence-only models (TemBERTure, DeepTM)")
    sb.add_bullet(tf, "First GNN Tm regression on 100K AlphaFold structures")
    sb.add_bullet(tf, "Validates multi-modal approach (PLM + 3D structure)")

    # ════════════════════════════════════════════
    # SLIDE 12: Base Paper 1
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "BASE PAPER 1", "11")
    sb.add_footer(slide)
    
    # Subtitle
    tf = sb.add_text_box(slide, BODY_LEFT, 580000, BODY_W, 500000, 'Text Sub')
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '"ThermoAGT-GA: Restricted-Attention GNN for Thermal Stability Prediction"'
    run.font.name = FONT_NAME
    run.font.size = SECTION_HEADER_FONT
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Journal of Theoretical Biology, 2024  |  Predicts \u0394\u0394G using dual augmented gated-attention GNN"
    run.font.name = FONT_NAME
    run.font.size = BODY_SMALL_FONT
    run.font.color.rgb = GRAY_TEXT
    
    # Left: Key contributions
    tf = sb.add_text_box(slide, BODY_LEFT, 1150000, 4300000, 3200000, 'TextLeft')
    sb.add_section(tf, "KEY CONTRIBUTIONS")
    sb.add_bullet(tf, "Residue subgraphs (10-residue window) + contact-map edges")
    sb.add_bullet(tf, "Dual augmented gated-attention GNN (WT vs mutant)")
    sb.add_bullet(tf, "True GNN on protein graphs, handles single + multi mutants")
    sb.add_bullet(tf, "Strong performance vs prior stability tools")
    sb.add_blank_line(tf)
    sb.add_section(tf, "LIMITATIONS")
    sb.add_bullet(tf, "Predicts \u0394\u0394G, not absolute Tm")
    sb.add_bullet(tf, "Requires structure/contact maps")
    sb.add_bullet(tf, "Local subgraph focus, not full-protein")
    
    # Right: Our advancement
    tf = sb.add_text_box(slide, 4660000, 1150000, 4300000, 3200000, 'TextRight')
    sb.add_section(tf, "OUR ADVANCEMENT")
    sb.add_bullet(tf, "We predict absolute Tm, not relative \u0394\u0394G")
    sb.add_bullet(tf, "Full protein graphs, not 10-residue subgraphs")
    sb.add_bullet(tf, "100K scale training vs limited mutation datasets")
    sb.add_bullet(tf, "ESM-2 embeddings vs basic residue features")
    sb.add_bullet(tf, "Direct regression on diverse protein families")

    # ════════════════════════════════════════════
    # SLIDE 13: Base Paper 2
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "BASE PAPER 2", "12")
    sb.add_footer(slide)
    
    tf = sb.add_text_box(slide, BODY_LEFT, 580000, BODY_W, 500000, 'Text Sub')
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '"TemBERTure: Advancing Protein Thermostability Prediction with Deep Learning"'
    run.font.name = FONT_NAME
    run.font.size = SECTION_HEADER_FONT
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Bioinformatics Advances, 2024  |  BERT-based PLM for Tm prediction from sequences"
    run.font.name = FONT_NAME
    run.font.size = BODY_SMALL_FONT
    run.font.color.rgb = GRAY_TEXT
    
    tf = sb.add_text_box(slide, BODY_LEFT, 1150000, 4300000, 3200000, 'TextLeft')
    sb.add_section(tf, "KEY CONTRIBUTIONS")
    sb.add_bullet(tf, "Direct Tm prediction using BERT-like PLM embeddings")
    sb.add_bullet(tf, "Curated Tm dataset (TemBERTure-DB)")
    sb.add_bullet(tf, "Both classifier (thermophilic vs mesophilic) and regressor")
    sb.add_bullet(tf, "Good baseline for sequence-only approach")
    sb.add_blank_line(tf)
    sb.add_section(tf, "LIMITATIONS")
    sb.add_bullet(tf, "No 3D structure or graph information used")
    sb.add_bullet(tf, "Dataset bias (species, Tm distribution)")
    sb.add_bullet(tf, "Limited to Meltome-derived data")
    
    tf = sb.add_text_box(slide, 4660000, 1150000, 4300000, 3200000, 'TextRight')
    sb.add_section(tf, "OUR ADVANCEMENT")
    sb.add_bullet(tf, "We add 3D structural graphs on top of PLM embeddings")
    sb.add_bullet(tf, "GNN captures spatial contacts that sequences cannot encode")
    sb.add_bullet(tf, "100K AlphaFold structures vs Meltome-only data")
    sb.add_bullet(tf, "Multi-modal approach: ESM-2 + 3D structure")
    sb.add_bullet(tf, "GIN architecture for more expressive graph learning")
    sb.add_bullet(tf, "Higher correlation through structure-aware features")

    # ════════════════════════════════════════════
    # SLIDE 14: Conclusion & Future Scope
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "CONCLUSION & FUTURE SCOPE", "13")
    sb.add_footer(slide)
    
    tf = sb.add_text_box(slide, BODY_LEFT, BODY_TOP, 4300000, 3700000, 'TextLeft')
    sb.add_section(tf, "ACHIEVEMENTS")
    sb.add_bullet(tf, "End-to-end pipeline: FASTA \u2192 AlphaFold \u2192 Graph \u2192 Tm prediction")
    sb.add_bullet(tf, "Successfully trained on 100,000 AlphaFold structures")
    sb.add_bullet(tf, "Hybrid GIN model achieved Pearson r = 0.744")
    sb.add_bullet(tf, "Integrated ESM-2 PLM embeddings + 3D structural contacts")
    sb.add_bullet(tf, "Addressed research gap: first GNN + PLM for absolute Tm")
    sb.add_blank_line(tf)
    sb.add_section(tf, "KEY CONCLUSION")
    sb.add_bullet(tf, "3D structural information significantly improves Tm prediction over sequence-only methods")
    sb.add_bullet(tf, "GIN's injective aggregation captures structural motifs better than GCN")
    
    tf = sb.add_text_box(slide, 4660000, BODY_TOP, 4300000, 3700000, 'TextRight')
    sb.add_section(tf, "FUTURE SCOPE")
    sb.add_bullet(tf, "Scale training to full 943K dataset using cloud GPU cluster")
    sb.add_bullet(tf, "Explore Graph Attention Networks (GAT) for attention-weighted edges")
    sb.add_bullet(tf, "Add Graph Transformer architecture for global context")
    sb.add_bullet(tf, "Build inference web app for instant Tm prediction from PDB files")
    sb.add_bullet(tf, "Feature importance / saliency analysis for interpretability")
    sb.add_blank_line(tf)
    sb.add_section(tf, "LIMITATIONS")
    sb.add_bullet(tf, "16GB RAM limits local training scale")
    sb.add_bullet(tf, "AlphaFold structures are predictions, not experimental")
    sb.add_bullet(tf, "Single-condition Tm (no pH/buffer normalization)")

    # ════════════════════════════════════════════
    # SLIDE 15: References
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "REFERENCES", "14")
    sb.add_footer(slide)
    
    tf = sb.add_text_box(slide, BODY_LEFT, BODY_TOP, BODY_W, BODY_H, 'Text 2')
    
    refs = [
        ("ThermoAGT-GA (2024)", "Predicting protein thermal stability changes upon single and multi-point mutations via restricted attention subgraph neural network. Journal of Theoretical Biology, 2024."),
        ("ProSTAGE (2024)", "Predicting Effects of Mutations on Protein Stability by Fusing Structure and Sequence Embedding. Journal of Chemical Information and Modeling, 2024."),
        ("Stability Oracle (2024)", "Structure-based graph-transformer framework for identifying thermodynamically stable protein sequences. Nature Communications, 2024."),
        ("TemBERTure (2024)", "Advancing protein thermostability prediction with deep learning. Bioinformatics Advances, 2024."),
        ("DeepTM (2023)", "Deep learning algorithm for prediction of melting temperature of thermophilic proteins directly from sequences. Computational and Structural Biotechnology Journal, 2023."),
        ("Dataset", "TemStaPro-Major-30-imbal-training.fasta — Protein thermostability dataset (Bioinformatics, 2024)."),
        ("ESM-2", "Lin et al. (2023). Evolutionary-scale prediction of atomic-level protein structure with a language model. Science, 379."),
    ]
    
    for title, desc in refs:
        sb.add_section(tf, title)
        sb.add_bullet(tf, desc, Pt(10))
        sb.add_blank_line(tf)

    # ════════════════════════════════════════════
    # SLIDE 16: Thank You
    # ════════════════════════════════════════════
    slide = sb.add_slide()
    sb.add_header_bar(slide, "", None)
    sb.add_footer(slide)
    
    # Thank You text centered
    tf = sb.add_text_box(slide, Emu(914400), Emu(1400000), Emu(7315200), Emu(800000), 'Text 1b')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.name = FONT_NAME
    run.font.size = THANK_YOU_FONT
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT
    
    tf = sb.add_text_box(slide, Emu(914400), Emu(2300000), Emu(7315200), Emu(700000), 'Text 2b')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Questions & Discussion"
    run.font.name = FONT_NAME
    run.font.size = SUBTITLE_FONT
    run.font.bold = True
    run.font.color.rgb = MAROON
    
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Team-22  |  IBS-2 Project  |  Semester 4"
    run.font.name = FONT_NAME
    run.font.size = SUBTITLE_SMALL_FONT
    run.font.color.rgb = GRAY_TEXT
    
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Structure-Guided GNN for Protein Melting Temperature Prediction"
    run.font.name = FONT_NAME
    run.font.size = BODY_FONT
    run.font.color.rgb = GRAY_TEXT
    
    # ──────────────────────────────────────────
    # SAVE
    # ──────────────────────────────────────────
    prs.save(OUTPUT_PATH)
    print(f"DONE! Saved presentation to: {OUTPUT_PATH}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    build_presentation()
